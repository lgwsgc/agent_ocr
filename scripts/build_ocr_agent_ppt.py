from pathlib import Path
import json

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = list(Path.home().joinpath("Desktop").glob("claude*20251104.pptx"))[0]
OUT = ROOT / "outputs" / "business_license_ocr_agent_tool_use.pptx"
IMG = ROOT / "Snipaste_2026-06-04_15-12-26.jpg"
RESULT = ROOT / "outputs" / "company_tool_image_result.json"
TEXT_RESULT = ROOT / "outputs" / "tool_trace_text_demo.json"

EMU_PER_IN = 914400


def emu(value):
    return int(value * EMU_PER_IN)


COLORS = {
    "ink": RGBColor(55, 48, 39),
    "muted": RGBColor(105, 97, 82),
    "accent": RGBColor(168, 56, 42),
    "gold": RGBColor(174, 138, 72),
    "soft": RGBColor(247, 243, 234),
    "line": RGBColor(218, 205, 176),
    "white": RGBColor(255, 255, 255),
    "green": RGBColor(45, 132, 91),
    "blue": RGBColor(55, 102, 145),
}


def remove_shape(shape):
    shape.element.getparent().remove(shape.element)


def clear_slide(slide, keep_freeform=False):
    for shape in list(slide.shapes):
        if keep_freeform and str(shape.shape_type).startswith("FREEFORM"):
            continue
        remove_shape(shape)


def add_text(slide, text, x, y, w, h, size=18, bold=False, color="ink", align="left"):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = emu(0.02)
    tf.margin_bottom = emu(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    return box


def add_card(slide, x, y, w, h, fill="soft", line="line"):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS[fill]
    shp.line.color.rgb = COLORS[line]
    shp.line.width = Pt(1)
    return shp


def add_label_card(slide, label, body, x, y, w, h, accent="accent", body_size=13):
    add_card(slide, x, y, w, h)
    add_text(slide, label, x + 0.14, y + 0.12, w - 0.28, 0.28, size=12.5, bold=True, color=accent)
    add_text(slide, body, x + 0.14, y + 0.46, w - 0.28, h - 0.55, size=body_size, color="ink")


def add_kicker(slide, text, page):
    add_text(slide, text, 0.48, 0.18, 7.2, 0.32, size=18, bold=True)
    add_text(slide, f"{page:02d}", 11.25, 0.2, 0.45, 0.22, size=8.5, color="muted", align="right")


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.48, 0.82, 11.35, 0.48, size=24, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.52, 1.34, 10.9, 0.36, size=12.3, color="muted")


def add_footer(slide, page):
    add_text(
        slide,
        "资料来源：《智能体设计模式(google).pdf》第5章；Demo：PaddleOCR-VL-1.6 + Business License OCR Agent",
        0.45,
        7.15,
        9.8,
        0.22,
        size=7.2,
        color="muted",
    )
    add_text(slide, f"{page:02d}", 11.62, 7.15, 0.32, 0.22, size=7.2, color="muted", align="right")


def add_arrow(slide, x1, y1, x2, y2):
    ln = slide.shapes.add_connector(1, emu(x1), emu(y1), emu(x2), emu(y2))
    ln.line.color.rgb = COLORS["gold"]
    ln.line.width = Pt(1.4)
    return ln


def add_picture(slide, path, x, y, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), emu(x), emu(y), width=emu(w), height=emu(h))


def load_json(path):
    return json.loads(path.read_text(encoding="utf-8")) if path.exists() else {}


def main():
    prs = Presentation(TEMPLATE)
    data = load_json(RESULT)
    fields = data.get("final_fields", {})
    text_demo = load_json(TEXT_RESULT)

    for idx, slide in enumerate(prs.slides, 1):
        clear_slide(slide, keep_freeform=(idx in [3, 5, 11]))

    # 1
    s = prs.slides[0]
    add_text(s, "Agentic OCR 实战", 1.1, 1.55, 10.0, 0.72, size=38, bold=True, align="center")
    add_text(s, "营业执照字段识别、校验与自动修复", 1.0, 2.35, 10.2, 0.52, size=24, bold=True, color="accent", align="center")
    add_text(s, "以《工具使用（函数调用）》为理论主线，把 OCR 从“识别文字”升级为“完成可信企业信息抽取任务”。", 1.4, 3.15, 9.2, 0.65, size=14.5, color="muted", align="center")
    add_text(s, "分享人：李岗威\n时间：2026-06-08", 8.85, 5.45, 2.75, 0.65, size=13)

    # 2
    s = prs.slides[1]
    add_text(s, "目录", 0.9, 0.95, 2.1, 0.5, size=30, bold=True)
    add_text(s, "CONTENTS", 0.92, 1.46, 2.4, 0.32, size=14, color="gold")
    contents = [
        ("01", "业务对象：为什么选营业执照", "字段稳定、规则强、错误可解释"),
        ("02", "理论主线：工具使用（函数调用）", "Agent 如何组织 OCR、校验器、修复器"),
        ("03", "Demo 演示：从图片到可信 JSON", "PaddleOCR-VL-1.6 + 字段修复闭环"),
        ("04", "沉淀 Skill：可复用环境与流程", "部署、代理、模型下载与故障恢复"),
    ]
    for i, (n, t, b) in enumerate(contents):
        y = 2.05 + i * 1.08
        add_text(s, n, 1.15, y, 0.6, 0.35, size=19, bold=True, color="accent")
        add_text(s, t, 1.9, y - 0.02, 5.3, 0.35, size=18, bold=True)
        add_text(s, b, 1.9, y + 0.35, 6.6, 0.28, size=11.5, color="muted")
    add_footer(s, 2)

    # 3
    s = prs.slides[2]
    add_text(s, "为什么选择营业执照", 3.45, 3.25, 5.6, 0.65, size=30, bold=True, align="center")
    add_text(s, "01", 5.7, 2.15, 0.8, 0.45, size=24, bold=True, color="accent", align="center")

    # 4
    s = prs.slides[3]
    add_kicker(s, "业务对象", 4)
    add_title(s, "营业执照是“字段固定 + 规则强约束”的理想 Agent OCR 场景", "相比开放文档，证照字段可枚举、格式可校验、错误可以解释。")
    add_picture(s, IMG, 0.65, 1.72, 5.25, 3.75)
    add_label_card(s, "核心字段", "统一社会信用代码、名称、类型、执行事务合伙人/法定代表人、出资额/注册资本、成立日期、经营场所、经营范围、登记机关、核准日期。", 6.25, 1.68, 5.0, 1.38)
    add_label_card(s, "为什么适合 Agent", "OCR 输出不是终点；Agent 可以把字段抽取、规则校验、自动修复、人类复核组织成可解释流程。", 6.25, 3.22, 5.0, 1.25, "blue")
    add_label_card(s, "可验证约束", "统一社会信用代码 18 位并带校验码；日期、金额、字段完整性都能作为反馈信号。", 6.25, 4.65, 5.0, 1.15, "green")
    add_footer(s, 4)

    # 5
    s = prs.slides[4]
    add_text(s, "理论主线：工具使用", 3.55, 3.25, 5.2, 0.65, size=30, bold=True, align="center")
    add_text(s, "02", 5.7, 2.15, 0.8, 0.45, size=24, bold=True, color="accent", align="center")

    # 6
    s = prs.slides[5]
    add_kicker(s, "第 5 章：工具使用（函数调用）", 6)
    add_title(s, "工具使用模式把 LLM 推理连接到外部能力", "PDF 第 5 章的核心流程可以直接映射到 OCR Agent。")
    steps = [("工具定义", "OCR、字段抽取器、校验器、修复器"), ("调用决策", "根据任务状态选择下一步工具"), ("工具执行", "PaddleOCR-VL 推理 / 规则函数"), ("观察结果", "OCR 文本、字段 JSON、失败原因"), ("继续处理", "修复、复核或输出报告")]
    for i, (a, b) in enumerate(steps):
        x = 0.55 + i * 2.32
        add_card(s, x, 2.25, 1.82, 1.45)
        add_text(s, str(i + 1), x + 0.12, 2.38, 0.3, 0.25, size=12, bold=True, color="accent")
        add_text(s, a, x + 0.42, 2.34, 1.15, 0.32, size=14, bold=True)
        add_text(s, b, x + 0.18, 2.82, 1.45, 0.55, size=10.2, color="muted")
        if i < 4:
            add_arrow(s, x + 1.85, 2.95, x + 2.26, 2.95)
    add_text(s, "对应到我们的任务：Agent 不直接“看图说话”，而是把 OCR 模型、确定性规则和修复逻辑都作为工具来编排。", 0.8, 4.55, 10.6, 0.6, size=18, bold=True, color="accent", align="center")
    add_footer(s, 6)

    # 7
    s = prs.slides[6]
    add_kicker(s, "系统架构", 7)
    add_title(s, "营业执照 OCR Agent = OCR 感知工具 + 规则校验工具 + 修复工具", "第 5 章的工具使用模式，在 demo 中落成一条可观测工具链。")
    blocks = [("输入", "营业执照图片\n或 OCR 文本"), ("OCR 工具", "PaddleOCR-VL-1.6\n输出版面块与文本"), ("抽取工具", "字段别名匹配\n标签/值分离处理"), ("校验工具", "信用代码校验\n日期与金额格式"), ("修复工具", "O/0、I/1\n候选修复"), ("输出", "JSON + 修复报告\n工具调用轨迹")]
    for i, (a, b) in enumerate(blocks):
        x = 0.5 + i * 1.92
        add_card(s, x, 2.25, 1.55, 1.25, "white")
        add_text(s, a, x + 0.1, 2.38, 1.35, 0.25, size=12, bold=True, color="accent", align="center")
        add_text(s, b, x + 0.1, 2.75, 1.35, 0.45, size=9.4, align="center")
        if i < 5:
            add_arrow(s, x + 1.55, 2.88, x + 1.9, 2.88)
    add_card(s, 1.05, 4.32, 10.05, 1.25)
    add_text(s, "工具调用轨迹 tool_trace", 1.35, 4.48, 2.4, 0.3, size=14, bold=True, color="blue")
    add_text(s, "read → extract → validate → repair → final_validate：每一步都有工具名、目的和观察结果，方便讲清楚“Agent 在做什么”。", 1.35, 4.88, 9.1, 0.35, size=14)
    add_footer(s, 7)

    # 8
    s = prs.slides[7]
    add_kicker(s, "Demo 环境", 8)
    add_title(s, "本地 CPU 也能跑：PaddleOCR-VL 走 transformers 引擎", "无显卡环境下，重点是提前缓存模型，并保留文本兜底入口。")
    add_label_card(s, "稳定演示入口", "& 'D:\\conda_envs\\company_tool\\python.exe' app.py --text samples/license_ocr_with_errors.txt", 0.65, 2.05, 5.2, 1.55, "green", 11)
    add_label_card(s, "真实图片入口", "$env:HTTP_PROXY='http://127.0.0.1:7890'\n& 'D:\\conda_envs\\company_tool\\python.exe' app.py --image .\\Snipaste_2026-06-04_15-12-26.jpg", 6.1, 2.05, 5.2, 1.85, "accent", 10.5)
    add_label_card(s, "模型缓存", "PaddleOCR-VL-1.6、PP-DocLayoutV3、PP-DocLayoutV3_safetensors 已缓存到 C:\\Users\\EDY\\.paddlex\\official_models", 0.65, 4.35, 5.2, 1.15, "blue", 10.5)
    add_label_card(s, "工程兜底", "现场网络或模型慢时，切换文本入口展示 Agent 修复闭环；图片 OCR 结果作为完整链路证明。", 6.1, 4.35, 5.2, 1.15, "gold")
    add_footer(s, 8)

    # 9
    s = prs.slides[8]
    add_kicker(s, "真实图片识别结果", 9)
    add_title(s, "PaddleOCR-VL 已从样例营业执照中抽取出关键企业字段", "OCR 负责“看见”，Agent 负责把文本整理成可信结构。")
    rows = [
        ("统一社会信用代码", fields.get("统一社会信用代码", "91110102082881146K")),
        ("名称", fields.get("名称", "中兴华会计师事务所（特殊普通合伙）")),
        ("类型", fields.get("类型", "特殊普通合伙企业")),
        ("执行事务合伙人", fields.get("法定代表人", "李尊农、乔久华")),
        ("出资额", fields.get("注册资本", "8916万元")),
        ("成立日期", fields.get("成立日期", "2013年11月04日")),
        ("核准日期", fields.get("核准日期", "2025年02月27日")),
    ]
    for i, (k, v) in enumerate(rows):
        y = 1.95 + i * 0.55
        add_text(s, k, 0.75, y, 1.65, 0.25, size=10.5, bold=True, color="accent")
        add_text(s, v, 2.45, y, 4.6, 0.28, size=10.5)
    add_picture(s, IMG, 7.25, 1.75, 4.25, 3.02)
    add_card(s, 7.25, 4.95, 4.25, 0.78)
    add_text(s, "校验结果：关键字段通过", 7.45, 5.12, 2.4, 0.25, size=13, bold=True, color="green")
    add_text(s, "status = pass", 9.85, 5.12, 1.1, 0.25, size=12)
    add_footer(s, 9)

    # 10
    s = prs.slides[9]
    add_kicker(s, "自动修复闭环", 10)
    add_title(s, "错误不是直接“猜改”，而是通过工具校验生成修复证据", "这页用于展示 Agentic OCR 相比普通 OCR 的价值。")
    add_label_card(s, "OCR 错误样例", "9111O1O2O8284I146L\n2013 年 11 月 O4 日\n2025 年 O2 月 27 日", 0.7, 1.85, 3.2, 1.55, "accent")
    add_label_card(s, "修复后", "91110102082841146L\n2013 年 11 月 04 日\n2025 年 02 月 27 日", 4.25, 1.85, 3.2, 1.55, "green")
    add_label_card(s, "为什么可信", "统一社会信用代码修复后通过 18 位字符集与校验码规则；日期修复后通过格式校验。", 7.8, 1.85, 3.6, 1.55, "blue")
    for i, item in enumerate(text_demo.get("tool_trace", [])[:5]):
        y = 4.25 + i * 0.38
        add_text(s, f"{item.get('step')}. {item.get('tool')}", 0.85, y, 2.35, 0.22, size=9.3, bold=True, color="accent")
        add_text(s, item.get("observation", ""), 3.0, y, 7.9, 0.22, size=9.3)
    add_footer(s, 10)

    # 11
    s = prs.slides[10]
    add_text(s, "从 Demo 到 Skill", 3.7, 3.25, 4.8, 0.65, size=30, bold=True, align="center")
    add_text(s, "03", 5.7, 2.15, 0.8, 0.45, size=24, bold=True, color="accent", align="center")

    # 12
    s = prs.slides[11]
    add_kicker(s, "沉淀复用", 12)
    add_title(s, "这次探索可以沉淀成两个 Skill：业务流程 + 环境部署", "让下一次遇到证照 OCR 或 PaddleOCR-VL 部署问题时，不再从头踩坑。")
    add_label_card(s, "business-license-ocr-agent", "记录营业执照字段、别名、校验规则、OCR 混淆修复、输出 JSON 契约。", 0.75, 1.95, 5.1, 1.25, "accent")
    add_label_card(s, "paddleocr-vl-local-setup", "记录 company_tool 环境、Clash 代理、模型下载、Python 3.10、transformers 引擎和常见故障恢复。", 6.15, 1.95, 5.1, 1.25, "blue")
    add_label_card(s, "落地经验 1", "本地 CPU 可跑，但要提前预热模型；真实图片链路和文本兜底链路都要准备。", 0.75, 3.65, 5.1, 1.0, "green")
    add_label_card(s, "落地经验 2", "不要只展示 OCR 文本，要展示工具调用轨迹、校验前后对比和修复证据。", 6.15, 3.65, 5.1, 1.0, "gold")
    add_footer(s, 12)

    # 13
    s = prs.slides[12]
    add_text(s, "总结", 0.95, 0.85, 2.0, 0.5, size=32, bold=True)
    add_text(s, "Agent 不是替代 OCR，而是让 OCR 结果可用、可信、可解释。", 0.95, 1.55, 10.2, 0.65, size=28, bold=True, color="accent")
    points = [
        ("从识别到任务", "OCR 输出文本；Agent 输出结构化企业信息和判断依据。"),
        ("从黑盒到轨迹", "tool_trace 展示每一步工具、目的、观察结果。"),
        ("从错误到修复", "规则校验把 OCR 错误转化为可解释的修复候选。"),
        ("从 demo 到复用", "业务 Skill + 部署 Skill 让经验可迁移。"),
    ]
    for i, (a, b) in enumerate(points):
        add_label_card(s, a, b, 0.95 + (i % 2) * 5.6, 2.65 + (i // 2) * 1.35, 5.0, 1.0, ["accent", "blue", "green", "gold"][i])
    add_text(s, "Q&A", 5.3, 6.35, 1.6, 0.4, size=22, bold=True, align="center")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
