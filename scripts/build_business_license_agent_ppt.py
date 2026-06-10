from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = Path("D:/OCR_Research/company_share/营业执照AgentOCR技术分享_v5.pptx")
LICENSE_IMAGE = ROOT / "data" / "188.jpg"


W, H = 13.333, 7.5


COLORS = {
    "white": RGBColor(255, 255, 255),
    "ink": RGBColor(48, 52, 59),
    "muted": RGBColor(102, 112, 123),
    "light": RGBColor(244, 248, 251),
    "line": RGBColor(218, 228, 236),
    "blue": RGBColor(15, 94, 142),
    "cyan": RGBColor(27, 158, 207),
    "green": RGBColor(69, 176, 94),
    "red": RGBColor(190, 57, 42),
    "gold": RGBColor(173, 128, 45),
}


def inch(value: float):
    return Inches(value)


def set_fill(shape, color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color]


def set_line(shape, color: str = "line", width: float = 0.8):
    shape.line.color.rgb = COLORS[color]
    shape.line.width = Pt(width)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    bold: bool = False,
    color: str = "ink",
    align=PP_ALIGN.LEFT,
    font: str = "Microsoft YaHei",
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.03)
    tf.margin_right = inch(0.03)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    return box


def add_para(slide, lines, x, y, w, h, *, size=15, color="ink"):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.03)
    tf.margin_right = inch(0.03)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(8)
    return box


def add_card(slide, x, y, w, h, *, fill="white", line="line", radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    set_fill(shape, fill)
    set_line(shape, line)
    return shape


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 1.0, 0.48, 10.8, 0.55, size=26, bold=True)
    if subtitle:
        add_text(slide, subtitle, 1.02, 1.08, 10.8, 0.35, size=12.5, color="muted")


def add_brand_bg(slide, page=None):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inch(W), inch(H))
    set_fill(bg, "white")
    bg.line.fill.background()

    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inch(0.22), inch(H))
    set_fill(left, "blue")
    left.line.fill.background()

    tri1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, inch(0.96), inch(3.9))
    set_fill(tri1, "green")
    tri1.line.fill.background()
    tri1.rotation = 180

    tri2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, 0, inch(3.25), inch(1.14), inch(4.25))
    set_fill(tri2, "cyan")
    tri2.line.fill.background()

    add_text(slide, "上海致宇信息技术有限公司", 9.25, 7.13, 2.8, 0.2, size=7.5, color="muted", align=PP_ALIGN.RIGHT)
    add_text(slide, "用技术简化知识工作", 1.0, 7.13, 2.0, 0.2, size=7.5, color="muted")
    if page is not None:
        add_text(slide, f"{page:02d}", 12.02, 7.12, 0.35, 0.2, size=7.5, color="muted", align=PP_ALIGN.RIGHT)


def add_section_slide(prs, number, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide)
    add_text(slide, number, 5.55, 2.13, 0.9, 0.44, size=24, bold=True, color="red", align=PP_ALIGN.CENTER)
    add_text(slide, title, 2.35, 3.12, 8.5, 0.65, size=31, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, 2.55, 3.86, 8.1, 0.38, size=13.5, color="muted", align=PP_ALIGN.CENTER)
    return slide


def add_metric(slide, label, value, x, y, color="blue"):
    add_card(slide, x, y, 2.45, 0.86, fill="light", line="line")
    add_text(slide, value, x + 0.16, y + 0.12, 2.1, 0.28, size=19, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.12, y + 0.51, 2.15, 0.2, size=8.8, color="muted", align=PP_ALIGN.CENTER)


def add_chip(slide, text, x, y, color="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(1.55), inch(0.35))
    set_fill(shape, color)
    shape.line.fill.background()
    add_text(slide, text, x + 0.06, y + 0.07, 1.42, 0.16, size=9, bold=True, color="white", align=PP_ALIGN.CENTER)


def add_flow(slide, items, x, y, w, h):
    gap = 0.16
    box_w = (w - gap * (len(items) - 1)) / len(items)
    for i, item in enumerate(items):
        bx = x + i * (box_w + gap)
        add_card(slide, bx, y, box_w, h, fill="light")
        add_text(slide, str(i + 1), bx + 0.08, y + 0.11, 0.22, 0.18, size=8.5, bold=True, color="red")
        add_text(slide, item[0], bx + 0.34, y + 0.08, box_w - 0.46, 0.25, size=11.5, bold=True)
        add_text(slide, item[1], bx + 0.13, y + 0.45, box_w - 0.26, h - 0.55, size=9.1, color="muted", align=PP_ALIGN.CENTER)


def add_grid_flow(slide, items, x, y, w, h, *, cols=3):
    gap_x = 0.22
    gap_y = 0.22
    rows = (len(items) + cols - 1) // cols
    box_w = (w - gap_x * (cols - 1)) / cols
    box_h = (h - gap_y * (rows - 1)) / rows
    for i, item in enumerate(items):
        row = i // cols
        col = i % cols
        bx = x + col * (box_w + gap_x)
        by = y + row * (box_h + gap_y)
        add_card(slide, bx, by, box_w, box_h, fill="light")
        add_text(slide, str(i + 1), bx + 0.12, by + 0.14, 0.22, 0.18, size=8.5, bold=True, color="red")
        add_text(slide, item[0], bx + 0.43, by + 0.12, box_w - 0.58, 0.22, size=12.5, bold=True)
        add_text(slide, item[1], bx + 0.18, by + 0.48, box_w - 0.36, box_h - 0.58, size=9.5, color="muted", align=PP_ALIGN.CENTER)


def build_deck():
    prs = Presentation()
    prs.slide_width = inch(W)
    prs.slide_height = inch(H)

    # 1 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide)
    add_text(slide, "营业执照 Agent OCR 实践", 1.15, 1.62, 10.5, 0.65, size=36, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "从 PaddleOCR-VL 识别到可信字段 JSON", 1.2, 2.42, 10.4, 0.46, size=22, bold=True, color="red", align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "工具调用 / LangChain / 大模型语义后处理 / Skill 沉淀",
        1.8,
        3.12,
        9.1,
        0.34,
        size=14.5,
        color="muted",
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "分享人：李岗威", 8.95, 5.55, 2.5, 0.25, size=12.5)
    add_text(slide, "时间：2026-06-10", 8.95, 5.92, 2.5, 0.25, size=12.5)

    # 2 目录
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 2)
    add_text(slide, "目录", 1.05, 0.88, 1.5, 0.46, size=30, bold=True)
    add_text(slide, "CONTENTS", 1.08, 1.35, 1.65, 0.22, size=10.5, color="gold")
    contents = [
        ("01", "营业执照常识：制度、版面与字段变化", "为什么它适合作为 Agent OCR 的业务对象"),
        ("02", "OCR 难点：从看见文字到可信字段", "长文本、地址、字段错配和语义型错误"),
        ("03", "理论主线：工具使用与 LangChain", "用函数调用把 OCR、校验、修复组织成流程"),
        ("04", "Demo 实践：从图片到可信 JSON", "PaddleOCR-VL-1.6 + 大模型语义后处理"),
        ("05", "沉淀与复盘：可复用流程和 Skill", "环境、代理、模型下载、默认参数和故障恢复"),
    ]
    for i, (num, title, sub) in enumerate(contents):
        y = 1.95 + i * 0.92
        add_text(slide, num, 1.18, y, 0.5, 0.3, size=18, bold=True, color="red")
        add_text(slide, title, 1.82, y - 0.01, 6.4, 0.3, size=17, bold=True)
        add_text(slide, sub, 1.84, y + 0.34, 7.1, 0.22, size=10.5, color="muted")

    # 3
    add_section_slide(prs, "01", "营业执照常识", "先理解证照制度和字段变化，再谈识别难点")

    # 4 制度变化
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 4)
    add_title(slide, "营业执照的发展：从“证件图片”到“企业身份入口”", "营业执照承载的是主体资格、经营边界和监管入口，因此识别结果天然需要可信。")
    timeline = [
        ("传统注册号", "各地工商登记规则不统一，证照字段和编号规则差异较大。"),
        ("三证合一 / 一照一码", "统一社会信用代码成为企业身份主键。"),
        ("新版横版执照", "版面更稳定，二维码、公示系统入口和字段布局更标准。"),
        ("字段持续调整", "法定代表人、经营者、出资额、经营场所等随主体类型变化。"),
    ]
    for i, (title, body) in enumerate(timeline):
        x = 1.08 + i * 2.72
        add_card(slide, x, 2.12, 2.35, 2.25, fill="light")
        add_text(slide, f"0{i + 1}", x + 0.13, 2.27, 0.38, 0.22, size=10.5, bold=True, color="red")
        add_text(slide, title, x + 0.14, 2.72, 2.0, 0.28, size=15, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.18, 3.22, 1.96, 0.68, size=10.8, color="muted", align=PP_ALIGN.CENTER)
    add_text(slide, "分享视角：营业执照不是“随便 OCR 一下”的图片，而是有制度约束、字段语义和校验规则的业务对象。", 1.2, 5.35, 10.2, 0.42, size=19, bold=True, color="red", align=PP_ALIGN.CENTER)

    # 5 字段变化
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 5)
    add_title(slide, "版面和字段会变：同一位置不一定代表同一业务含义", "营业执照识别难点不止是字准，还包括主体类型驱动的字段解释。")
    rows = [
        ("公司", "法定代表人", "注册资本", "住所", "经营范围"),
        ("合伙企业", "执行事务合伙人", "出资额", "主要经营场所", "经营范围"),
        ("个体工商户", "经营者", "资金数额/无", "经营场所", "经营范围"),
        ("分支机构", "负责人", "无注册资本", "营业场所", "经营范围"),
    ]
    x0, y0 = 1.06, 1.86
    widths = [1.55, 2.25, 1.65, 2.35, 2.45]
    headers = ["主体类型", "负责人字段", "资本字段", "地址字段", "经营边界"]
    for col, header in enumerate(headers):
        add_card(slide, x0 + sum(widths[:col]), y0, widths[col], 0.46, fill="blue", line="blue", radius=False)
        add_text(slide, header, x0 + sum(widths[:col]) + 0.04, y0 + 0.12, widths[col] - 0.08, 0.16, size=9.5, bold=True, color="white", align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 0.46 + r * 0.58
        for col, val in enumerate(row):
            fill = "light" if r % 2 == 0 else "white"
            add_card(slide, x0 + sum(widths[:col]), y, widths[col], 0.58, fill=fill, radius=False)
            add_text(slide, val, x0 + sum(widths[:col]) + 0.05, y + 0.17, widths[col] - 0.1, 0.18, size=9.6, align=PP_ALIGN.CENTER)
    add_para(
        slide,
        [
            "这也是为什么 demo 不能只做普通 OCR：同一张图中可能有正文、说明、二维码旁注、印章和脚注。",
            "Agent 的价值在于把“字段抽取、规则校验、上下文修复、输出解释”拆成可观测工具链。",
        ],
        1.08,
        5.0,
        10.5,
        0.76,
        size=15.5,
    )

    # 6
    add_section_slide(prs, "02", "OCR 难点", "营业执照真正难的，是把识别结果变成可信字段")

    # 7 难点
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 7)
    add_title(slide, "从图片到字段，中间有四类典型错误", "这些错误正好可以作为 Agent OCR 的实战样本。")
    cards = [
        ("版面噪声", "国徽、印章、二维码、边框、水印会干扰版面块定位。", "green"),
        ("字段错配", "把旧执照或说明文字中的内容误认为目标字段。", "red"),
        ("长文本截断", "经营范围很长，换行和括号会影响完整性。", "blue"),
        ("语义型错字", "如“临颍/临颖”“造价/浩价”“记账/记帐”，单靠格式校验不够。", "gold"),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 1.05 + (i % 2) * 5.45
        y = 1.86 + (i // 2) * 1.62
        add_card(slide, x, y, 4.95, 1.18, fill="light")
        add_chip(slide, title, x + 0.18, y + 0.17, color)
        add_text(slide, body, x + 0.24, y + 0.62, 4.42, 0.32, size=13, color="ink")
    add_text(slide, "结论：OCR 模型负责“看见”，但业务系统还需要“理解、校验、修复、解释”。", 1.1, 5.7, 10.3, 0.34, size=20, bold=True, color="red", align=PP_ALIGN.CENTER)

    # 8
    add_section_slide(prs, "03", "工具使用与 LangChain", "把知识点落到 demo 的架构里")

    # 9 工具调用
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 9)
    add_title(slide, "工具使用：让模型调用外部能力，而不是只生成文本", "对应 PDF 中《工具使用（函数调用）》章节：定义工具、调用工具、观察结果、决定下一步。")
    add_flow(
        slide,
        [
            ("工具定义", "OCR、字段抽取、校验、规则修复、语义修复"),
            ("调用决策", "根据当前任务状态选择下一步工具"),
            ("执行工具", "本地模型、Python 函数或远程大模型"),
            ("观察结果", "文本、JSON、校验失败原因、耗时"),
            ("继续处理", "再次校验、修复或输出报告"),
        ],
        0.92,
        2.15,
        10.95,
        1.28,
    )
    add_card(slide, 1.2, 4.35, 10.1, 1.1, fill="white")
    add_text(slide, "映射到本次 demo", 1.45, 4.54, 1.65, 0.24, size=13, bold=True, color="blue")
    add_text(slide, "LangChain StructuredTool 把每个动作封装成工具；Agent 记录 tool_trace，最后输出字段 JSON、修复动作和运行摘要。", 3.05, 4.54, 7.75, 0.32, size=14.5)

    # 10 LangChain
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 10)
    add_title(slide, "LangChain 在 demo 中承担“工具编排层”", "它不是为了炫框架，而是让每个步骤可替换、可追踪、可复用。")
    flow = [
        ("输入工具", "读取图片或 OCR 文本"),
        ("OCR 工具", "PaddleOCR-VL-1.6"),
        ("抽取工具", "中文字段结构化"),
        ("校验工具", "信用代码 / 日期 / 金额"),
        ("修复工具", "规则修复 + LLM 语义修复"),
        ("输出工具", "JSON + 摘要 + 路径"),
    ]
    add_flow(slide, flow, 0.78, 1.95, 11.45, 1.5)
    add_para(
        slide,
        [
            "设计原则：可确定的错误先用规则修；需要上下文和语义的错误再交给大模型。",
            "这能避免把所有问题都丢给 LLM，也方便后续替换 OCR 模型或大模型供应商。",
        ],
        1.05,
        4.62,
        10.4,
        0.75,
        size=16,
    )

    # 11
    add_section_slide(prs, "04", "Demo 实践", "从营业执照图片到可信字段 JSON")

    # 12 Demo pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 12)
    add_title(slide, "Demo 流程：图片进入，可信 JSON 出来", "当前默认参数已经指向样例营业执照图片，运行 app.py 即可走完整链路。")
    if LICENSE_IMAGE.exists():
        slide.shapes.add_picture(str(LICENSE_IMAGE), inch(1.0), inch(1.68), width=inch(4.75))
    add_grid_flow(
        slide,
        [
            ("图片输入", "data/188.jpg"),
            ("OCR", "PaddleOCR-VL-1.6"),
            ("字段抽取", "中文字段"),
            ("校验", "格式和规则"),
            ("LLM 修复", "语义后处理"),
            ("结果保存", "outputs/*.json"),
        ],
        6.1,
        1.86,
        5.4,
        2.38,
        cols=3,
    )
    add_text(slide, "演示命令", 6.18, 5.03, 1.1, 0.24, size=13, bold=True, color="blue")
    add_card(slide, 6.12, 5.38, 5.38, 0.62, fill="light")
    add_text(slide, "python app.py", 6.35, 5.58, 4.8, 0.18, size=15, bold=True)

    # 13 运行摘要
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 13)
    add_title(slide, "这次运行确实经过了大模型后处理", "摘要中能看到 OCR 模型、LLM 模型、耗时和结果保存路径。")
    add_metric(slide, "OCR 模型", "PaddleOCR-VL-1.6", 1.05, 1.88, "blue")
    add_metric(slide, "大模型", "claude-sonnet-4-6", 3.83, 1.88, "red")
    add_metric(slide, "大模型耗时", "14.256 秒", 6.61, 1.88, "green")
    add_metric(slide, "自动修复字段", "7 个", 9.39, 1.88, "gold")
    add_card(slide, 1.08, 3.52, 10.8, 1.46, fill="light")
    add_text(slide, "输出路径", 1.35, 3.74, 1.2, 0.24, size=13, bold=True, color="blue")
    add_text(slide, r"D:\OCR_Research\agent\outputs\license_image_agent_result.json", 2.35, 3.74, 8.8, 0.24, size=14.5, bold=True)
    add_text(slide, "完整图片链路在 CPU 上耗时较长，主要时间花在本地 PaddleOCR-VL 推理；LLM 后处理只占十几秒。", 1.35, 4.24, 9.8, 0.28, size=13, color="muted")

    # 14 字段结果
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 14)
    add_title(slide, "最终输出字段：用中文字段名，方便业务同学直接理解", "字段不用英文包装，demo 直接输出营业执照语义字段。")
    final_rows = [
        ("统一社会信用代码", "91110108306767909G"),
        ("名称", "北京索伯尼国际科技有限公司"),
        ("类型", "有限责任公司（自然人独资）"),
        ("法定代表人", "刘福喜"),
        ("注册资本", "2000万元"),
        ("成立日期", "2014年08月25日"),
        ("营业期限", "2014年08月25日至2034年08月24日"),
        ("住所", "北京市海淀区中关村大街49号9号楼B座四层405室"),
        ("核准日期", "2019年11月04日"),
    ]
    x0, y0 = 1.04, 1.7
    for i, (k, v) in enumerate(final_rows):
        y = y0 + i * 0.46
        fill = "light" if i % 2 == 0 else "white"
        add_card(slide, x0, y, 2.15, 0.42, fill=fill, radius=False)
        add_card(slide, x0 + 2.15, y, 8.25, 0.42, fill=fill, radius=False)
        add_text(slide, k, x0 + 0.1, y + 0.12, 1.85, 0.14, size=8.9, bold=True, color="blue", align=PP_ALIGN.CENTER)
        add_text(slide, v, x0 + 2.3, y + 0.12, 7.8, 0.14, size=9.2)
    add_text(slide, "经营范围字段较长，在 JSON 中完整保留；PPT 中只展示关键短字段。", 1.06, 6.15, 10.0, 0.24, size=12.5, color="muted")

    # 15 修复动作
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 15)
    add_title(slide, "大模型后处理负责“语义修复”，但必须保留证据链", "本次图片识别中，LLM 根据上下文修复了 7 个字段。")
    repairs = [
        ("统一社会信用代码", "92440300MA5FPG0XXH", "91110108306767909G"),
        ("类型", "个体工商户", "有限责任公司（自然人独资）"),
        ("法定代表人", "薛丁川", "刘福喜"),
        ("成立日期", "2019年07月12日", "2014年08月25日"),
        ("住所", "深圳市宝安区...佛商大厦027", "北京市海淀区中关村大街...405室"),
    ]
    headers = ["字段", "修复前", "修复后"]
    widths = [2.05, 4.05, 4.45]
    x0, y0 = 1.0, 1.83
    for c, htxt in enumerate(headers):
        add_card(slide, x0 + sum(widths[:c]), y0, widths[c], 0.42, fill="blue", line="blue", radius=False)
        add_text(slide, htxt, x0 + sum(widths[:c]) + 0.06, y0 + 0.11, widths[c] - 0.12, 0.14, size=9.2, bold=True, color="white", align=PP_ALIGN.CENTER)
    for r, row in enumerate(repairs):
        y = y0 + 0.42 + r * 0.6
        for c, val in enumerate(row):
            fill = "light" if r % 2 == 0 else "white"
            add_card(slide, x0 + sum(widths[:c]), y, widths[c], 0.6, fill=fill, radius=False)
            add_text(slide, val, x0 + sum(widths[:c]) + 0.08, y + 0.18, widths[c] - 0.16, 0.16, size=8.7, align=PP_ALIGN.CENTER)
    add_text(slide, "注意：大模型不能替代权威库。地址、省市区、主体名称等字段，后续最好接入行政区划库和企业信用公示数据做二次校验。", 1.0, 5.75, 10.8, 0.42, size=14.5, color="red", bold=True, align=PP_ALIGN.CENTER)

    # 16 模型边界
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 16)
    add_title(slide, "PaddleOCR-VL 是视觉文档理解模型，大模型承担语义后处理", "这页专门回答分享时最容易被问到的边界问题。")
    add_label = [
        ("PaddleOCR-VL-1.6", "多模态文档解析/视觉 OCR 模型，擅长版面、文本、表格和图片区域解析。它可以理解文档结构，但不是通用对话大模型。", "blue"),
        ("LLM 后处理", "利用上下文、字段语义和常识修复 OCR 错字、字段错配、经营范围截断等问题。", "red"),
        ("规则与权威库", "信用代码、日期、金额可规则校验；省市区地址建议接入行政区划库，避免只靠大模型猜。", "green"),
    ]
    for i, (title, body, color) in enumerate(add_label):
        add_card(slide, 1.05, 1.85 + i * 1.24, 10.55, 0.92, fill="light")
        add_chip(slide, title, 1.28, 2.08 + i * 1.24, color)
        add_text(slide, body, 3.05, 2.05 + i * 1.24, 8.15, 0.26, size=13)
    add_text(slide, "我的判断：让 OCR/VL 负责感知，让规则负责确定性，让 LLM 负责语义候选，再用业务库兜底。", 1.2, 5.75, 10.2, 0.36, size=19, bold=True, color="red", align=PP_ALIGN.CENTER)

    # 17
    add_section_slide(prs, "05", "沉淀与复盘", "把一次 demo 变成可复用的方法")

    # 18 总结
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bg(slide, 18)
    add_title(slide, "沉淀下来的是一套 Agent OCR 方法，而不只是一段脚本", "可复用的东西，才是这次分享真正有价值的部分。")
    takeaways = [
        ("默认参数", "app.py 默认走图片链路，输出路径固定，现场不需要临时拼命输参数。"),
        ("运行摘要", "只打印模型、耗时、是否调用 LLM、保存位置，避免现场被大段 JSON 淹没。"),
        ("环境 Skill", "记录 conda、代理、SSL_CERT_FILE、模型缓存和常见问题恢复。"),
        ("业务 Skill", "记录营业执照字段、别名、校验规则、修复策略和讲解口径。"),
    ]
    for i, (title, body) in enumerate(takeaways):
        x = 1.05 + (i % 2) * 5.45
        y = 1.9 + (i // 2) * 1.35
        add_card(slide, x, y, 4.95, 0.96, fill="light")
        add_text(slide, title, x + 0.18, y + 0.15, 1.3, 0.24, size=13, bold=True, color=["blue", "red", "green", "gold"][i])
        add_text(slide, body, x + 1.45, y + 0.15, 3.25, 0.36, size=11.5)
    add_text(slide, "Q&A", 5.18, 5.82, 1.8, 0.45, size=25, bold=True, align=PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build_deck())
